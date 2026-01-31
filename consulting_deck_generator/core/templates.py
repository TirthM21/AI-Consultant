from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import base64
from typing import Dict, Any, List, Optional
from .models import SlideContent, SlideData
from .chart_generator import ChartGenerator

class SlideTemplates:
    """McKinsey-style slide templates"""
    
    def __init__(self):
        self.mckinsey_colors = {
            'primary': RGBColor(0, 164, 228),
            'secondary': RGBColor(0, 38, 58),
            'accent': RGBColor(255, 107, 53),
            'gray': RGBColor(139, 105, 151),
            'text': RGBColor(0, 0, 0),
            'background': RGBColor(255, 255, 255)
        }
    
    def create_title_slide(self, prs, title: str, subtitle: str = None, client: str = None):
        """Create McKinsey title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_font = title_shape.text_frame.paragraphs[0].font
        title_font.bold = True
        title_font.size = Pt(44)
        title_font.color.rgb = self.mckinsey_colors['secondary']
        
        # Add subtitle if provided
        if subtitle or client:
            subtitle_text = f"{client or ''}\\n{subtitle or ''}"
            left = Inches(1)
            top = Inches(3.5)
            width = Inches(8)
            height = Inches(1)
            
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text.strip()
            subtitle_font = subtitle_frame.paragraphs[0].font
            subtitle_font.size = Pt(18)
            subtitle_font.color.rgb = self.mckinsey_colors['primary']
            subtitle_font.italic = True
        
        # Add McKinsey-style footer
        self._add_footer(slide, "Proprietary & Confidential")
        
        return slide
    
    def create_two_column_slide(self, prs, title: str, left_content: List[str], right_content: List[str]):
        """Create two-column layout slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        slide.shapes.title.text = title
        title_font = slide.shapes.title.text_frame.paragraphs[0].font
        title_font.bold = True
        title_font.size = Pt(36)
        title_font.color.rgb = self.mckinsey_colors['secondary']
        
        # Remove existing content placeholder
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.element.getparent().remove(shape.element)
        
        # Create two columns
        col_width = Inches(4)
        col_height = Inches(4)
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(1), Inches(2), col_width, col_height)
        left_frame = left_box.text_frame
        left_frame.text = "\\n".join(left_content)
        self._style_bullet_text(left_frame)
        
        # Right column  
        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), col_width, col_height)
        right_frame = right_box.text_frame
        right_frame.text = "\\n".join(right_content)
        self._style_bullet_text(right_frame)
        
        self._add_footer(slide)
        
        return slide
    
    def create_chart_slide(self, prs, title: str, chart_base64: str, caption: str = None):
        """Create slide with chart"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        slide.shapes.title.text = title
        title_font = slide.shapes.title.text_frame.paragraphs[0].font
        title_font.bold = True
        title_font.size = Pt(36)
        title_font.color.rgb = self.mckinsey_colors['secondary']
        
        # Remove content placeholder
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.element.getparent().remove(shape.element)
        
        # Add chart image
        chart_bytes = base64.b64decode(chart_base64)
        chart_image = io.BytesIO(chart_bytes)
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        slide.shapes.add_picture(chart_image, left, top, width, height)
        
        # Add caption if provided
        if caption:
            caption_top = Inches(6.8)
            caption_height = Inches(0.5)
            
            caption_box = slide.shapes.add_textbox(Inches(1), caption_top, width, caption_height)
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            
            caption_font = caption_frame.paragraphs[0].font
            caption_font.size = Pt(14)
            caption_font.color.rgb = self.mckinsey_colors['primary']
            caption_font.italic = True
        
        self._add_footer(slide)
        
        return slide
    
    def create_framework_slide(self, prs, title: str, framework_data: List[List[str]]):
        """Create framework/2x2 matrix slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        slide.shapes.title.text = title
        title_font = slide.shapes.title.text_frame.paragraphs[0].font
        title_font.bold = True
        title_font.size = Pt(36)
        title_font.color.rgb = self.mckinsey_colors['secondary']
        
        # Remove content placeholder
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.element.getparent().remove(shape.element)
        
        # Create 2x2 framework
        cell_width = Inches(3.5)
        cell_height = Inches(2)
        
        for i, row in enumerate(framework_data):
            for j, cell_content in enumerate(row):
                left = Inches(1 + j * 4)
                top = Inches(2 + i * 2.5)
                
                # Create background box
                box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, cell_width, cell_height
                )
                box.fill.solid()
                box.fill.fore_color.rgb = self.mckinsey_colors['light_gray'] if i == 0 and j == 0 else RGBColor(255, 255, 255)
                box.line.color.rgb = self.mckinsey_colors['gray']
                box.line.width = Pt(1)
                
                # Add text
                text_frame = box.text_frame
                text_frame.text = cell_content
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                
                cell_font = text_frame.paragraphs[0].font
                cell_font.size = Pt(12)
                if i == 0 and j == 0:  # Highlight key cell
                    cell_font.color.rgb = RGBColor(255, 255, 255)
                    cell_font.bold = True
                else:
                    cell_font.color.rgb = self.mckinsey_colors['text']
        
        self._add_footer(slide)
        
        return slide
    
    def _style_bullet_text(self, text_frame):
        """Apply McKinsey bullet point styling"""
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self.mckinsey_colors['text']
            
            # Add bullets for non-empty paragraphs
            if paragraph.text.strip():
                paragraph.level = 0
                # Use default bullet style
    
    def _add_footer(self, slide, text: str = "Proprietary & Confidential"):
        """Add McKinsey-style footer"""
        left = Inches(0.5)
        top = Inches(7)
        width = Inches(9)
        height = Inches(0.3)
        
        footer_box = slide.shapes.add_textbox(left, top, width, height)
        footer_frame = footer_box.text_frame
        
        # Date on left
        footer_frame.text = f"{datetime.now().strftime('%B %d, %Y')} | {text}"
        
        footer_font = footer_frame.paragraphs[0].font
        footer_font.size = Pt(8)
        footer_font.color.rgb = self.mckinsey_colors['gray']
        footer_font.italic = True